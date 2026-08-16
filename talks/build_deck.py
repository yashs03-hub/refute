"""Builds talks/reagent_judging_5min.pptx.

Structure, per direct instruction 2026-08-16: the pipeline diagram (question
-> layer 1 -> HANDOFF -> layer 2), then a worked example, then a layer-3
close on why this beats a frontier model alone. The worked example is real,
not illustrative - the influenza/fibroblast-necroptosis case a layer-1
session actually produced, run through THIS repo's real `refute intake`
live in this same session (the near-tie bleomycin_lung/apoptosis_resistance
result and the honest "outside what any protocol can model" outcome are the
actual CLI output, not a mockup).

Regenerate: python talks/build_deck.py
"""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# --- palette -----------------------------------------------------------------

BG = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1A, 0x1D, 0x21)
MUTED = RGBColor(0x5B, 0x63, 0x6C)
NAVY = RGBColor(0x0F, 0x2A, 0x4A)
ACCENT = RGBColor(0xC0, 0x3B, 0x2E)      # the number-that-matters red
ACCENT_GREEN = RGBColor(0x1E, 0x7B, 0x4D)
RULE = RGBColor(0xDD, 0xE1, 0xE5)
CODE_BG = RGBColor(0x14, 0x17, 0x1A)
CODE_TEXT = RGBColor(0xE8, 0xEA, 0xED)
CODE_MUTED = RGBColor(0x8A, 0x92, 0x9C)
BOX_L1 = RGBColor(0xEC, 0xF1, 0xF7)
BOX_L2 = RGBColor(0xFB, 0xEF, 0xEC)

SANS = "Arial"
MONO = "Courier New"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _no_autosize(tf):
    el = tf._txBody
    bodyPr = el.find(qn("a:bodyPr"))
    for tag in ("a:normAutofit", "a:spAutoFit"):
        existing = bodyPr.find(qn(tag))
        if existing is not None:
            bodyPr.remove(existing)
    bodyPr.append(el.makeelement(qn("a:noAutofit"), {}))


def _box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    _no_autosize(tf)
    return tb, tf


def _set(p, text, size, color=INK, bold=False, font=SANS, italic=False):
    r = p.add_run()
    r.text = text if text else " "
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    return r


def _add(p, text, size, color=INK, bold=False, font=SANS, italic=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    return r


def _kicker(slide, text, y=Inches(0.55)):
    _, tf = _box(slide, Inches(0.7), y, Inches(11.9), Inches(0.4))
    p = tf.paragraphs[0]
    _set(p, text.upper(), 13, MUTED, bold=True)
    for r in p.runs:
        r.font._rPr.set("spc", "150")


def _title(slide, text, y=Inches(0.95), size=32):
    _, tf = _box(slide, Inches(0.7), y, Inches(11.9), Inches(1.1))
    p = tf.paragraphs[0]
    _set(p, text, size, NAVY, bold=True)


def _rule(slide, y):
    ln = slide.shapes.add_connector(1, Inches(0.7), y, Inches(12.63), y)
    ln.line.color.rgb = RULE
    ln.line.width = Pt(1)


def _footer(slide, n):
    _, tf = _box(slide, Inches(0.7), Inches(7.08), Inches(6), Inches(0.35))
    p = tf.paragraphs[0]
    _set(p, "refute — re:AGENT, Track B", 10, MUTED)
    _, tf2 = _box(slide, Inches(11.9), Inches(7.08), Inches(0.7), Inches(0.35))
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    _set(p2, str(n), 10, MUTED)


def _code_block(slide, x, y, w, h, lines, size=15):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.adjustments[0] = 0.04
    box.fill.solid()
    box.fill.fore_color.rgb = CODE_BG
    box.line.fill.background()
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.3)
    tf.margin_top = Inches(0.22)
    tf.margin_bottom = Inches(0.22)
    _no_autosize(tf)
    for i, (text, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        _set(p, text, size, color, bold=bold, font=MONO)
    return box


def _bullets(slide, x, y, w, h, items, size=17, gap=10, color=INK, lead_color=NAVY):
    _, tf = _box(slide, x, y, w, h)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        if isinstance(item, tuple):
            lead, rest = item
            _set(p, f"{lead}  ", size, lead_color, bold=True)
            _add(p, rest, size, color)
        else:
            _set(p, item, size, color)
    return tf


def _big_stat(slide, x, y, w, number, label, number_color=ACCENT, number_size=44):
    _, tf = _box(slide, x, y, w, Inches(1.0))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _set(p, number, number_size, number_color, bold=True)
    _, tf2 = _box(slide, x, y + Inches(0.85), w, Inches(0.7))
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    p2.word_wrap = True
    _set(p2, label, 13, MUTED)


def _pipeline_box(slide, x, y, w, h, title, lines, fill, title_size=17):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.adjustments[0] = 0.05
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = NAVY
    box.line.width = Pt(1.25)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.35)
    tf.margin_right = Inches(0.35)
    tf.margin_top = Inches(0.2)
    tf.margin_bottom = Inches(0.2)
    _no_autosize(tf)
    p = tf.paragraphs[0]
    _set(p, title, title_size, NAVY, bold=True)
    for line in lines:
        pp = tf.add_paragraph()
        pp.space_before = Pt(5)
        _set(pp, f"·  {line}", 13.5, INK)
    return box


def _arrow_label(slide, y, label, size=12.5):
    _, tf = _box(slide, Inches(0.7), y, Inches(1.4), Inches(0.5))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _set(p, "▼", 20, NAVY, bold=True)
    if label:
        _, tf2 = _box(slide, Inches(2.3), y + Inches(0.02), Inches(10.3), Inches(0.6))
        p2 = tf2.paragraphs[0]
        _set(p2, label, size, MUTED, italic=True)


# -------------------------------------------------------------------------

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

n = 0


def new_slide():
    global n
    n += 1
    s = prs.slides.add_slide(blank)
    _bg(s)
    return s


# --- 1. Title ------------------------------------------------------------

s = new_slide()
_, tf = _box(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(1.4))
p = tf.paragraphs[0]
_set(p, "refute", 56, NAVY, bold=True)
_, tf = _box(s, Inches(0.95), Inches(3.4), Inches(11.5), Inches(1.0))
p = tf.paragraphs[0]
_set(p, "Grounded hypothesis vetting, and the handoff to experimental design.",
     22, INK)
_, tf = _box(s, Inches(0.95), Inches(4.15), Inches(11.5), Inches(0.6))
p = tf.paragraphs[0]
_set(p, "Layer 2 of a two-layer system  ·  re:AGENT, Track B", 15, MUTED)
_rule(s, Inches(6.3))
_footer(s, n)

# --- 2. The pipeline ---------------------------------------------------

s = new_slide()
_kicker(s, "The architecture")
_title(s, "Two layers, one handoff.")

_, tf = _box(s, Inches(4.9), Inches(1.85), Inches(3.5), Inches(0.5))
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
_set(p, "a question", 16, INK, italic=True)

_arrow_label(s, Inches(2.35), "")

_pipeline_box(
    s, Inches(1.4), Inches(2.75), Inches(10.5), Inches(1.95),
    "LAYER 1 — due diligence, to exhaustion",
    [
        "check it against the computational data",
        "check it against the literature",
        "rule out what can be ruled out",
        "loop until nothing left changes the answer",
    ],
    BOX_L1,
)

_arrow_label(
    s, Inches(4.85),
    "HANDOFF — everything found, everything ruled out, and the residual "
    "the bench must settle",
)

_pipeline_box(
    s, Inches(1.4), Inches(5.65), Inches(10.5), Inches(1.55),
    "LAYER 2 — experimental design  (this repository)",
    [
        "pick the assay, size it",
        "judge whether it can be assessed, and at what fidelity",
    ],
    BOX_L2,
)
_footer(s, n)

# --- 3. Worked example: the hypothesis ----------------------------------

s = new_slide()
_kicker(s, "Worked example")
_title(s, "A scientist's hypothesis, run through the full pipeline.")
_, tf = _box(s, Inches(0.7), Inches(1.95), Inches(11.9), Inches(0.55))
p = tf.paragraphs[0]
_set(p, "Mouse whole-lung scRNA-seq, influenza A PR8, five timepoints, "
        "29,513 cells.", 16, MUTED)

_code_block(s, Inches(0.7), Inches(2.65), Inches(11.9), Inches(1.55), [
    ('"fibroblasts are dying during acute infection... to hyper recruit', CODE_TEXT, False),
    (' immune cells, particular the interferon responsive and the', CODE_TEXT, False),
    (' adventitial fibroblasts"', CODE_TEXT, False),
    ("  — message 185, as asked, typos included", CODE_MUTED, True),
], size=15)

_bullets(s, Inches(0.7), Inches(4.45), Inches(11.9), Inches(2.1), [
    "Refined at message 191 into three separable propositions: a "
    "TRANSITION (adventitial → interferon-responsive), a DEATH "
    "(inflammatory), and a RECRUITMENT CONSEQUENCE.",
    "“Nothing to prove that. That's my hypothesis, but we should just "
    "see where the data and the analysis take us.” — the scientist, "
    "same message.",
], size=16, gap=16)
_footer(s, n)

# --- 4. What layer 1 caught ----------------------------------------------

s = new_slide()
_kicker(s, "What layer 1 caught")
_title(s, "It contradicted its own earlier claim — and caught it.")
_bullets(s, Inches(0.7), Inches(2.0), Inches(6.6), Inches(3.4), [
    ("C7, message 185 —", "“the dataset contains no explicit "
     "influenza viral-feature gene.”"),
    ("C1, message 44, same agent —", "had already found it and quantified "
     "it: “Flu: 16.72% of 67,211 cells express it.”"),
    ("What caught it —", "not new data, not a better model. The system's "
     "own trace: “C7 directly contradicts C1 made by the same agent "
     "at message 44.”"),
], size=17, gap=16)
_, tf = _box(s, Inches(0.7), Inches(5.45), Inches(6.6), Inches(1.3))
p = tf.paragraphs[0]
_set(p, "The invalidating evidence sat in its own context for 141 "
        "messages. Nothing looked.", 17, ACCENT, bold=True)

_code_block(s, Inches(7.7), Inches(2.0), Inches(4.9), Inches(4.75), [
    ("13 claims scored", CODE_TEXT, True),
    ("", CODE_TEXT, False),
    ("1  contradicted  (C7)", ACCENT, True),
    ("7  abstained on", CODE_MUTED, False),
    ("   \"nothing relevant", CODE_MUTED, False),
    ("    in hand\"", CODE_MUTED, False),
    ("5  upheld", ACCENT_GREEN, False),
    ("", CODE_TEXT, False),
    ("negative control —", CODE_TEXT, True),
    ("a separate thread,", CODE_MUTED, False),
    ("8 claims, all correct,", CODE_MUTED, False),
    ("144 papers read,", CODE_MUTED, False),
    ("nothing killed", ACCENT_GREEN, True),
], size=14)
_footer(s, n)

# --- 5. The handoff --------------------------------------------------------

s = new_slide()
_kicker(s, "The handoff")
_title(s, "What crossed: findings, what's ruled out, one residual.")
_bullets(s, Inches(0.7), Inches(2.0), Inches(11.9), Inches(2.5), [
    ("Four experiments ruled out —", "simple acute depletion (counts never "
     "fall below baseline), mitochondrial fraction as a death readout "
     "(uniformly zero in this object), the 3.7-fold ISG number as quotable "
     "(direction survives, magnitude doesn't), any replicate-supported "
     "inference (one orig.ident)."),
    ("The residual, stated as the brief —", "adventitial fibroblasts "
     "acquire necroptosis-biased SIGNALLING POTENTIAL; whether that "
     "executes as death, and whether it drives recruitment, is "
     "unanswerable in this object."),
], size=16, gap=20)
_rule(s, Inches(4.75))
_, tf = _box(s, Inches(0.7), Inches(5.0), Inches(11.9), Inches(1.6))
p = tf.paragraphs[0]
_set(p, "“protein-level phospho-MLKL/RIPK3 and cleaved-caspase assays "
        "would be needed”", 19, NAVY, bold=True)
p2 = tf.add_paragraph()
p2.space_before = Pt(6)
_set(p2, "— the analysis names its own experiment. That sentence is what "
         "crosses to layer 2.", 15, MUTED, italic=True)
_footer(s, n)

# --- 6. Layer 2's real response ------------------------------------------

s = new_slide()
_kicker(s, "Layer 2 — this repository, live")
_title(s, "Ran it through refute intake. It refused to force a match.")
_code_block(s, Inches(0.7), Inches(2.0), Inches(11.9), Inches(3.35), [
    ("$ refute intake \"a fibroblast-restricted perturbation in mouse", CODE_MUTED, False),
    ("  lung across an influenza PR8 timecourse: conditional Ripk3 or", CODE_MUTED, False),
    ("  Mlkl deletion... phospho-MLKL/RIPK3 readout...\"", CODE_MUTED, False),
    ("", CODE_TEXT, False),
    ("1. bleomycin_lung        7.1   lung, mouse, day, readout", CODE_TEXT, False),
    ("2. apoptosis_resistance  6.9   surviving, timepoint, viability,", CODE_TEXT, False),
    ("                               fibroblast, readout", CODE_TEXT, False),
    ("", CODE_TEXT, False),
    ("The leader is not clear of the field. Acting on it alone", ACCENT, True),
    ("discards a live alternative, so read the ranking rather", ACCENT, True),
    ("than the first line.", ACCENT, True),
    ("", CODE_TEXT, False),
    ("ready for the gate: no", ACCENT, True),
], size=14.5)
_, tf = _box(s, Inches(0.7), Inches(5.6), Inches(11.9), Inches(1.2))
p = tf.paragraphs[0]
_set(p, "Neither scaffold is a real fit — an in vivo genetic perturbation "
        "crossed with an infection timecourse is its own apparatus. "
        "Correctly refused, not forced.", 16, INK)
_footer(s, n)

# --- 7. Layer 2 grounds it anyway ------------------------------------------

s = new_slide()
_kicker(s, "Nothing prebuilt fits — grounded anyway")
_title(s, "What's buildable, and what isn't.", size=28)
_bullets(s, Inches(0.7), Inches(1.9), Inches(11.9), Inches(4.85), [
    ("Cre driver — precedented, not adventitial-specific.",
     "Pdgfra-CreERT2 already used in a PR8-family infection study (Jones "
     "2024, Science) — true adventitial drivers (Pi16/Dpt) validated only "
     "in skin, never lung."),
    ("Floxed Ripk3/Mlkl — a real gap.", "Conditional alleles exist; none "
     "has ever been crossed to a fibroblast Cre, in any tissue. A "
     "from-scratch cross, not an ordering formality."),
    ("p-MLKL detection — the strongest-grounded piece.", "IHC precedented "
     "on PR8-infected lung itself (Wang 2019), optimized on infected lung "
     "(Kelepouras 2024). Flow on dissociated cells has ZERO precedent — "
     "likely the same dissociation-loss problem the scRNA-seq hit."),
    ("The hypothesis itself is a real, unanswered gap.", "IRFibs are "
     "named, real cells from an actual PR8 paper (Boyd 2020, Nature) — "
     "which never looks at whether they die. No fibroblast-specific "
     "effect size exists to power against."),
    ("Ferroptosis arm — the weakest-grounded piece.", "No lung-flow "
     "protocol found anywhere for BODIPY-C11. Needs its own pilot first."),
], size=13, gap=8)
_footer(s, n)

# --- 8. Layer 3 — beats a frontier model ------------------------------------

s = new_slide()
_kicker(s, "Layer 3")
_title(s, "What a frontier model alone cannot do here.")
_bullets(s, Inches(0.7), Inches(2.05), Inches(11.9), Inches(4.5), [
    ("No self-catch —", "a model asked this question once does not "
     "re-read its own message 44 against its own message 185. C7 stayed "
     "wrong for 141 messages because nothing was built to look back."),
    ("No honest refusal —", "asked to design the experiment, a frontier "
     "model proposes one. It does not say “neither scaffold is a "
     "real fit” — it has no scaffolds to check against, and no "
     "incentive to say no."),
    ("No sourced numbers —", "a Cre driver line, a floxed allele, an "
     "antibody clone — a frontier model will name something plausible. "
     "This names something published, quoted, full-text-verified, or "
     "says NOT_REPORTED with the query attached."),
], size=18, gap=22)
_rule(s, Inches(5.7))
_, tf = _box(s, Inches(0.7), Inches(5.95), Inches(11.9), Inches(1.0))
p = tf.paragraphs[0]
_set(p, "The gap was never idea generation. It's catching your own error, "
        "and refusing when you don't actually know.", 18, NAVY, bold=True)
_footer(s, n)

# --- 9. Close ----------------------------------------------------------------

s = new_slide()
_kicker(s, "Close")
_title(s, "Two independently built systems. One real handoff.")
_bullets(s, Inches(0.7), Inches(2.1), Inches(11.9), Inches(3.6), [
    "Nothing on this path was invented for the deck — the hypothesis, "
    "the self-correction, the abstentions, the residual, and the intake "
    "run are all real output from real systems, this session.",
    "Layer 1 ran to exhaustion and handed off exactly one residual. "
    "Layer 2 checked it honestly against what's calibrated, found "
    "nothing that fits, and said so instead of forcing a number.",
    "That refusal, followed by real literature grounding instead of a "
    "guess, is the whole pitch in one worked example.",
], size=18, gap=20)
_footer(s, n)

out_path = "talks/reagent_judging_5min.pptx"
prs.save(out_path)
print(f"wrote {out_path}, {n} slides")
